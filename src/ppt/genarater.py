import os

from pptx import Presentation

from model.power_point import PowerPoint
from utils import remove_all_slides


class PPTGenerator:
    def __init__(self):
        pass

    @staticmethod
    def generate(power_point: PowerPoint, template_path: str, output_file_path: str) -> None:
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file '{template_path}' does not exist.")

        presensation = Presentation(template_path)
        remove_all_slides(presensation)
        presensation.core_properties.title = power_point.title
        for slide in power_point.slides:
            if slide.layout >= len(presensation.slide_layouts):
                slide_layout = presensation.slide_layouts[0]
            else:
                slide_layout = presensation.slide_layouts[slide.layout]

            new_slide = presensation.slides.add_slide(slide_layout)

            if new_slide.shapes.title:
                new_slide.shapes.title.text = slide.content.title

            for shape in new_slide.shapes:
                if shape.has_text_frame and not shape == new_slide.shapes.title:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    for point in slide.content.bullet_points:
                        p = text_frame.add_paragraph()
                        p.text = point
                        p.level = 0
                    break

            if slide.content.image_path:
                image_full_path = os.path.join(os.getcwd(), slide.content.image_path)
                if os.path.exists(image_full_path):
                    for shape in new_slide.placeholders:
                        if shape.placeholder_format.type == 18:
                            shape.insert_picture(image_full_path)
                            break
        presensation.save(output_file_path)
        print(f"Presentation saved to '{output_file_path}'")
